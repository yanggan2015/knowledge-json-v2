#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
从 Markdown 教程生成以框图为主的精美 PPTX 与视频幻灯片。

用法:
  python3 generate_ppt_from_article.py articles/React/chapters/001-*.md -o ppt_output
"""

from __future__ import annotations

import argparse
import asyncio
import math
import re
import subprocess
import textwrap
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

FONT_PATH = "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"
# 默认竖屏 9:16，适合手机观看
W, H = 1080, 1920
IS_PORTRAIT = True
DEFAULT_VOICE = "zh-CN-XiaoxiaoNeural"  # 女声；可选 zh-CN-YunxiNeural 男声

# 推荐发音人预设（--voice-preset）
VOICE_PRESETS: dict[str, list[str]] = {
    "female": ["zh-CN-XiaoxiaoNeural", "zh-CN-XiaoyiNeural"],
    "male": ["zh-CN-YunxiNeural", "zh-CN-YunjianNeural", "zh-CN-YunyangNeural"],
    "duo": ["zh-CN-XiaoxiaoNeural", "zh-CN-YunxiNeural"],  # 女声+男声轮换
    "mix": [
        "zh-CN-XiaoxiaoNeural",
        "zh-CN-YunxiNeural",
        "zh-CN-XiaoyiNeural",
        "zh-CN-YunjianNeural",
    ],
    "dialect": [
        "zh-CN-liaoning-XiaobeiNeural",
        "zh-CN-shaanxi-XiaoniNeural",
    ],
}

VOICE_LABELS: dict[str, str] = {
    "zh-CN-XiaoxiaoNeural": "晓晓 · 温柔女声（默认）",
    "zh-CN-XiaoyiNeural": "晓伊 · 活泼女声",
    "zh-CN-YunxiNeural": "云希 · 沉稳男声",
    "zh-CN-YunjianNeural": "云健 · 新闻男声",
    "zh-CN-YunyangNeural": "云扬 · 专业男声",
    "zh-CN-YunxiaNeural": "云夏 · 少年男声",
    "zh-CN-liaoning-XiaobeiNeural": "晓北 · 东北女声",
    "zh-CN-shaanxi-XiaoniNeural": "晓妮 · 陕西女声",
    "zh-HK-HiuGaaiNeural": "粤语 · 女声",
    "zh-TW-HsiaoChenNeural": "台湾 · 女声",
}

# 深色科技风配色
C_BG_TOP = (12, 18, 35)
C_BG_BOT = (28, 38, 58)
C_TITLE = (248, 250, 252)
C_SUB = (148, 163, 184)
C_LINE = (71, 85, 105)
PALETTE = [
    (56, 189, 248),   # sky
    (52, 211, 153),   # emerald
    (167, 139, 250),  # violet
    (251, 191, 36),   # amber
    (244, 114, 182),  # pink
    (94, 234, 212),   # teal
]


@dataclass
class Node:
    label: str
    detail: str = ""
    color_idx: int = 0


@dataclass
class Slide:
    title: str
    kind: str  # title | flow | layers | cards | pipeline | hub | summary
    nodes: List[Node] = field(default_factory=list)
    edges: List[Tuple[int, int]] = field(default_factory=list)
    subtitle: str = ""
    caption: str = ""
    narration: str = ""


def set_orientation(portrait: bool) -> None:
    global W, H, IS_PORTRAIT
    IS_PORTRAIT = portrait
    W, H = (1080, 1920) if portrait else (1920, 1080)


def narration_for(slide: Slide) -> str:
    if slide.narration:
        return slide.narration
    if slide.kind == "title":
        sub = slide.subtitle.replace("｜", "，") if slide.subtitle else ""
        return f"欢迎学习本节内容：{slide.title}。{sub}我们将通过框图方式讲解核心知识。"
    parts = [f"接下来看{slide.title}。"]
    if slide.caption:
        parts.append(slide.caption)
    for node in slide.nodes[:6]:
        seg = node.label
        if node.detail:
            seg += f"，{node.detail}"
        parts.append(seg)
    text = "。".join(parts)
    return text[:600]


async def synthesize_narration(
    text: str,
    out: Path,
    voice: str,
    rate: str = "+0%",
    volume: str = "+0%",
    pitch: str = "+0Hz",
) -> None:
    import edge_tts
    communicate = edge_tts.Communicate(text, voice, rate=rate, volume=volume, pitch=pitch)
    await communicate.save(str(out))


def audio_duration(path: Path) -> float:
    r = subprocess.run(
        [
            "ffprobe", "-v", "error", "-show_entries", "format=duration",
            "-of", "default=noprint_wrappers=1:nokey=1", str(path),
        ],
        capture_output=True,
        text=True,
        check=True,
    )
    return float(r.stdout.strip())


def resolve_voices(
    voice: str,
    voices_arg: list[str] | None,
    preset: str | None,
) -> list[str]:
    if preset:
        key = preset.lower()
        if key not in VOICE_PRESETS:
            raise SystemExit(f"未知预设: {preset}，可选: {', '.join(VOICE_PRESETS)}")
        return list(VOICE_PRESETS[key])
    if voices_arg:
        return voices_arg
    return [voice]


def pick_voice_for_slide(
    index: int,
    slide: Slide,
    voices: list[str],
    mode: str,
) -> str:
    if len(voices) == 1:
        return voices[0]
    if mode == "single":
        return voices[0]
    if mode == "by_kind":
        if slide.kind == "title":
            return voices[0]
        if slide.kind in ("flow", "pipeline", "layers"):
            return voices[1 % len(voices)]
        return voices[index % len(voices)]
    # rotate（默认）：按页轮换
    return voices[index % len(voices)]


def generate_all_narrations(
    slides: List[Slide],
    audio_dir: Path,
    voices: list[str],
    rate: str = "+0%",
    volume: str = "+0%",
    pitch: str = "+0Hz",
    voice_mode: str = "rotate",
) -> List[float]:
    async def _run() -> List[float]:
        audio_dir.mkdir(parents=True, exist_ok=True)
        durations: List[float] = []
        used: list[str] = []
        for i, slide in enumerate(slides, 1):
            text = narration_for(slide)
            mp3 = audio_dir / f"slide_{i:02d}.mp3"
            v = pick_voice_for_slide(i - 1, slide, voices, voice_mode)
            used.append(v)
            await synthesize_narration(text, mp3, v, rate=rate, volume=volume, pitch=pitch)
            durations.append(audio_duration(mp3))
        # 记录每页发音人
        (audio_dir / "voices_used.txt").write_text(
            "\n".join(f"slide_{j+1:02d}: {used[j]}" for j in range(len(used))),
            encoding="utf-8",
        )
        return durations

    return asyncio.run(_run())


def print_voice_catalog(locale_prefix: str = "zh") -> None:
    import edge_tts
    voices = asyncio.run(edge_tts.list_voices())
    print("=== 推荐预设 (--voice-preset) ===")
    for name, ids in VOICE_PRESETS.items():
        labels = [VOICE_LABELS.get(v, v) for v in ids]
        print(f"  {name}: {' | '.join(labels)}")
    print("\n=== 全部中文相关语音 (--voice / --voices) ===")
    for v in sorted(voices, key=lambda x: x["ShortName"]):
        if v["Locale"].startswith(locale_prefix):
            label = VOICE_LABELS.get(v["ShortName"], "")
            extra = f"  {label}" if label else ""
            print(f"  {v['ShortName']}  {v['Gender']}  {v['Locale']}{extra}")


def _margin() -> int:
    return 48 if IS_PORTRAIT else 72


def _clean(text: str) -> str:
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text.strip()


def _font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_PATH, size)


def _gradient_bg() -> Image.Image:
    img = Image.new("RGB", (W, H))
    for y in range(H):
        t = y / H
        r = int(C_BG_TOP[0] + (C_BG_BOT[0] - C_BG_TOP[0]) * t)
        g = int(C_BG_TOP[1] + (C_BG_BOT[1] - C_BG_TOP[1]) * t)
        b = int(C_BG_TOP[2] + (C_BG_BOT[2] - C_BG_TOP[2]) * t)
        ImageDraw.Draw(img).line([(0, y), (W, y)], fill=(r, g, b))
    return img


def _draw_glow_orbs(draw: ImageDraw.ImageDraw) -> None:
    if IS_PORTRAIT:
        orbs = [(180, 320, 140, (30, 58, 90)), (900, 500, 200, (40, 50, 80)), (540, 1500, 160, (35, 55, 75))]
    else:
        orbs = [(220, 200, 180, (30, 58, 90)), (1680, 320, 240, (40, 50, 80)), (1500, 820, 160, (35, 55, 75))]
    for cx, cy, rad, col in orbs:
        draw.ellipse([cx - rad, cy - rad, cx + rad, cy + rad], outline=col, width=2)


def _rounded_box(
    draw: ImageDraw.ImageDraw,
    xy: Tuple[int, int, int, int],
    fill: Tuple[int, int, int],
    outline: Tuple[int, int, int],
    radius: int = 20,
    width: int = 2,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def _text_centered(
    draw: ImageDraw.ImageDraw,
    box: Tuple[int, int, int, int],
    title: str,
    detail: str = "",
    title_font: ImageFont.FreeTypeFont | None = None,
    detail_font: ImageFont.FreeTypeFont | None = None,
    color: Tuple[int, int, int] = C_TITLE,
) -> None:
    title_font = title_font or _font(28)
    detail_font = detail_font or _font(22)
    x1, y1, x2, y2 = box
    tw = x2 - x1
    lines = textwrap.wrap(title, width=max(8, tw // 28))
    detail_lines = textwrap.wrap(detail, width=max(10, tw // 24))[:3] if detail else []
    total_h = len(lines) * 34 + len(detail_lines) * 28 + (10 if detail_lines else 0)
    y = y1 + (y2 - y1 - total_h) // 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=title_font)
        tw_line = bbox[2] - bbox[0]
        draw.text((x1 + (tw - tw_line) // 2, y), line, font=title_font, fill=color)
        y += 34
    if detail_lines:
        y += 6
        for line in detail_lines:
            bbox = draw.textbbox((0, 0), line, font=detail_font)
            tw_line = bbox[2] - bbox[0]
            draw.text((x1 + (tw - tw_line) // 2, y), line, font=detail_font, fill=C_SUB)
            y += 28


def _arrow(draw: ImageDraw.ImageDraw, x1: int, y1: int, x2: int, y2: int, color: Tuple[int, int, int]) -> None:
    draw.line([(x1, y1), (x2, y2)], fill=color, width=3)
    ang = math.atan2(y2 - y1, x2 - x1)
    sz = 12
    p1 = (x2 - sz * math.cos(ang - 0.4), y2 - sz * math.sin(ang - 0.4))
    p2 = (x2 - sz * math.cos(ang + 0.4), y2 - sz * math.sin(ang + 0.4))
    draw.polygon([(x2, y2), p1, p2], fill=color)


def _header(draw: ImageDraw.ImageDraw, title: str, subtitle: str = "", page: int = 0) -> int:
    m = _margin()
    draw.rectangle([0, 0, W, 6], fill=PALETTE[0])
    title_font = _font(40 if IS_PORTRAIT else 44)
    draw.text((m, 36), title[:48], font=title_font, fill=C_TITLE)
    bottom = 88
    if subtitle:
        draw.text((m, 88), subtitle[:100], font=_font(22), fill=C_SUB)
        bottom = 118
    if page:
        draw.text((W - m - 40, H - 48), f"{page:02d}", font=_font(22), fill=C_SUB)
    return bottom


def _content_area(header_bottom: int) -> Tuple[int, int, int, int]:
    """内容区矩形 (left, top, right, bottom)，尽量占满标题下方至页脚之间的空间。"""
    m = _margin()
    top = header_bottom + 16
    bottom = H - 56
    return (m, top, W - m, bottom)


def _layout_boxes_vertical(
    n: int,
    area: Tuple[int, int, int, int],
    min_box_h: int = 72,
) -> List[Tuple[int, int, int, int]]:
    left, top, right, bottom = area
    if n <= 0:
        return []
    height = bottom - top
    gap = max(10, height // (n * 6 + 4))
    total_gap = gap * max(0, n - 1)
    box_h = max(min_box_h, (height - total_gap) // n)
    if n * box_h + total_gap > height:
        box_h = (height - total_gap) // n
    boxes: List[Tuple[int, int, int, int]] = []
    y = top
    for _ in range(n):
        boxes.append((left, y, right, y + box_h))
        y += box_h + gap
    return boxes


def _layout_boxes_horizontal(
    n: int,
    area: Tuple[int, int, int, int],
    min_box_w: int = 100,
) -> List[Tuple[int, int, int, int]]:
    left, top, right, bottom = area
    if n <= 0:
        return []
    width = right - left
    height = bottom - top
    gap = max(14, width // (n * 8 + 2))
    total_gap = gap * max(0, n - 1)
    box_w = max(min_box_w, (width - total_gap) // n)
    if n * box_w + total_gap > width:
        box_w = (width - total_gap) // n
    boxes: List[Tuple[int, int, int, int]] = []
    x = left
    for _ in range(n):
        boxes.append((x, top, x + box_w, bottom))
        x += box_w + gap
    return boxes


def _layout_grid(
    n: int,
    area: Tuple[int, int, int, int],
    cols: int,
    min_box_h: int = 72,
) -> List[Tuple[int, int, int, int]]:
    left, top, right, bottom = area
    if n <= 0:
        return []
    cols = max(1, min(cols, n))
    rows = math.ceil(n / cols)
    width = right - left
    height = bottom - top
    gap_x = max(14, width // 16)
    gap_y = max(10, height // (rows * 6 + 2))
    total_gap_x = gap_x * (cols - 1)
    total_gap_y = gap_y * (rows - 1)
    box_w = (width - total_gap_x) // cols
    box_h = max(min_box_h, (height - total_gap_y) // rows)
    if rows * box_h + total_gap_y > height:
        box_h = (height - total_gap_y) // rows
    boxes: List[Tuple[int, int, int, int]] = []
    for i in range(n):
        row = i // cols
        col = i % cols
        items_in_row = min(cols, n - row * cols)
        row_w = items_in_row * box_w + (items_in_row - 1) * gap_x
        row_x = left + (width - row_w) // 2
        x = row_x + col * (box_w + gap_x)
        y = top + row * (box_h + gap_y)
        boxes.append((x, y, x + box_w, y + box_h))
    return boxes


def parse_mermaid(body: str) -> Tuple[List[Node], List[Tuple[int, int]], str]:
    """解析 flowchart / graph 为节点与边"""
    nodes_map: dict[str, int] = {}
    nodes: List[Node] = []
    edges: List[Tuple[int, int]] = []
    direction = "LR"
    if re.search(r"flowchart\s+TD|graph\s+TB|flowchart\s+TB", body):
        direction = "TB"

    for m in re.finditer(r"(\w+)\s*\[([^\]]+)\]", body):
        nid, label = m.group(1), _clean(m.group(2))
        if nid not in nodes_map:
            nodes_map[nid] = len(nodes)
            nodes.append(Node(label=label, color_idx=len(nodes) % len(PALETTE)))

    for m in re.finditer(r"(\w+)\s*-->\s*(\w+)", body):
        a, b = m.group(1), m.group(2)
        if a in nodes_map and b in nodes_map:
            edges.append((nodes_map[a], nodes_map[b]))

    # subgraph 标签作层名
    layers = re.findall(r"subgraph\s+([^\n{]+)", body)
    layer_hint = layers[0].strip() if layers else ""
    return nodes, edges, direction


SKIP_SECTIONS = frozenset({"延伸学习", "巩固建议", "常见误区与纠正", "最佳实践", "延伸阅读"})


def parse_article(path: Path) -> tuple[str, str, List[Slide]]:
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()
    title = _clean(lines[0][2:]) if lines and lines[0].startswith("# ") else path.stem
    meta = " ".join(_clean(l[2:]) for l in lines[1:8] if l.startswith("> "))

    slides: List[Slide] = [
        Slide(title=title, kind="title", subtitle=meta, caption="教程章节 · 框图导读")
    ]

    current = ""
    body_lines: List[str] = []

    def flush():
        nonlocal current, body_lines
        if not current:
            return
        if current in SKIP_SECTIONS:
            current, body_lines = "", []
            return
        body = "\n".join(body_lines).strip()
        if not body:
            current, body_lines = "", []
            return
        body = re.sub(r"^###\s+.+$", "", body, flags=re.MULTILINE)

        if current == "核心知识":
            concepts = re.findall(
                r"\*\*(\d+)\.\s*([^*]+)\*\*\s*\n+([^\n#*]+(?:\n(?!\*\*\d)[^\n#*]+)*)",
                body,
            )
            if concepts:
                nodes = [
                    Node(label=_clean(ct), detail=_clean(cb)[:80], color_idx=i)
                    for i, (_, ct, cb) in enumerate(concepts[:4])
                ]
                slides.append(Slide(title="核心知识图谱", kind="cards", nodes=nodes, caption=current))
                # 关系页：概念之间的逻辑链
                if len(nodes) >= 3:
                    slides.append(
                        Slide(
                            title="核心概念关系",
                            kind="flow",
                            nodes=nodes[:5],
                            edges=[(i, i + 1) for i in range(min(4, len(nodes) - 1))],
                            caption="概念递进关系",
                        )
                    )
            else:
                plain = _clean(re.sub(r"```[\s\S]*?```", "", body))
                slides.append(
                    Slide(
                        title=current,
                        kind="cards",
                        nodes=[Node(label=plain[:40], detail=plain[40:120], color_idx=0)],
                    )
                )
        elif "```mermaid" in body:
            mermaid = re.search(r"```mermaid\s*([\s\S]*?)```", body)
            m_body = mermaid.group(1) if mermaid else body
            nodes, edges, direction = parse_mermaid(m_body)
            if "subgraph" in m_body:
                slides.append(
                    Slide(title="架构分层", kind="layers", nodes=nodes, edges=edges, caption=current)
                )
            else:
                kind = "flow" if direction == "LR" else "pipeline"
                slides.append(
                    Slide(title="流程框图", kind=kind, nodes=nodes, edges=edges, caption=current)
                )
        else:
            numbered = [_clean(x) for x in re.findall(r"^\d+\.\s+(.+)$", body, re.MULTILINE)]
            if numbered:
                nodes = [Node(label=f"步骤 {i+1}", detail=t[:70], color_idx=i) for i, t in enumerate(numbered[:5])]
                slides.append(
                    Slide(
                        title=current,
                        kind="pipeline",
                        nodes=nodes,
                        edges=[(i, i + 1) for i in range(len(nodes) - 1)],
                    )
                )
            else:
                plain = _clean(re.sub(r"```[\s\S]*?```", "", body))
                chunks = textwrap.wrap(plain, width=36)[:4]
                nodes = [Node(label=c[:36], color_idx=i) for i, c in enumerate(chunks)]
                slides.append(Slide(title=current, kind="cards", nodes=nodes))

        current, body_lines = "", []

    for line in lines:
        if line.startswith("## "):
            flush()
            current = _clean(line[3:])
            body_lines = []
        elif line.startswith("---"):
            continue
        elif current:
            body_lines.append(line)
    flush()

    hub = Slide(
        title="本章要点回顾",
        kind="hub",
        nodes=[
            Node("概念模型", color_idx=0),
            Node("流程机制", color_idx=1),
            Node("工程实践", color_idx=2),
            Node("性能安全", color_idx=3),
            Node("持续学习", color_idx=4),
        ],
        caption="围绕本章主题的系统化掌握路径",
    )
    core = slides[0]
    body_slides = slides[1:]
    if len(body_slides) > 10:
        body_slides = body_slides[:10]
    slides = [core] + body_slides + [hub]
    return title, meta, slides


def _draw_node(draw: ImageDraw.ImageDraw, box: Tuple[int, int, int, int], node: Node) -> None:
    color = PALETTE[node.color_idx % len(PALETTE)]
    inner = (tuple(max(0, c - 80) for c in color))
    _rounded_box(draw, box, fill=(inner[0] // 4 + 20, inner[1] // 4 + 28, inner[2] // 4 + 45), outline=color, radius=18, width=3)
    _text_centered(draw, box, node.label, node.detail, _font(26), _font(20), C_TITLE)


def render_title(slide: Slide, img: Image.Image, page: int) -> None:
    draw = ImageDraw.Draw(img)
    _draw_glow_orbs(draw)
    draw.rectangle([0, 0, W, 6], fill=PALETTE[0])
    m = _margin()
    box_top = 140 if IS_PORTRAIT else 120
    box_bottom = H - 72
    _rounded_box(draw, (m, box_top, W - m, box_bottom), fill=(22, 32, 50), outline=PALETTE[0], radius=32, width=2)
    inner_h = box_bottom - box_top
    title_font = _font(52 if IS_PORTRAIT else 64)
    title_y = box_top + inner_h // 8
    draw.text((m + 40, title_y), slide.title, font=title_font, fill=C_TITLE)
    if slide.subtitle:
        draw.text((m + 40, title_y + 90), slide.subtitle[:80], font=_font(24), fill=C_SUB)
    if slide.caption:
        cap_y = title_y + 170
        _rounded_box(draw, (m + 40, cap_y, m + 420, cap_y + 56), fill=(56, 189, 248), outline=PALETTE[0], radius=12)
        draw.text((m + 56, cap_y + 16), slide.caption, font=_font(22), fill=C_TITLE)
    tags = re.findall(r"[^｜|]+", slide.subtitle or "")
    x = m + 40
    tag_y = box_top + inner_h // 2
    for i, tag in enumerate(tags[:6]):
        tag = tag.strip()
        if not tag:
            continue
        bw = min(W - 2 * m - 80, len(tag) * 20 + 36)
        _rounded_box(draw, (x, tag_y, x + bw, tag_y + 52), fill=(30, 45, 68), outline=PALETTE[(i + 1) % len(PALETTE)], radius=10)
        draw.text((x + 12, tag_y + 14), tag[:18], font=_font(20), fill=C_TITLE)
        x += bw + 16
        if x > W - m - 100:
            x = m + 40
            tag_y += 64
    draw.text((W - m - 20, H - 48), f"{page:02d}", font=_font(22), fill=C_SUB)


def render_flow(slide: Slide, img: Image.Image, page: int, vertical: bool = False) -> None:
    draw = ImageDraw.Draw(img)
    top = _header(draw, slide.title, slide.caption, page)
    nodes = slide.nodes[:6]
    n = len(nodes)
    if not n:
        return
    area = _content_area(top)
    use_vertical = vertical or IS_PORTRAIT
    if use_vertical:
        positions = _layout_boxes_vertical(n, area)
    else:
        positions = _layout_boxes_horizontal(n, area)

    for i, node in enumerate(nodes):
        _draw_node(draw, positions[i], node)

    edge_list = slide.edges or [(i, i + 1) for i in range(n - 1)]
    for a, b in edge_list:
        if a >= len(positions) or b >= len(positions):
            continue
        if use_vertical:
            cx = (positions[a][0] + positions[a][2]) // 2
            _arrow(draw, cx, positions[a][3], cx, positions[b][1], PALETTE[0])
        else:
            cy = (positions[a][1] + positions[a][3]) // 2
            _arrow(draw, positions[a][2], cy, positions[b][0], cy, PALETTE[0])


def render_cards(slide: Slide, img: Image.Image, page: int) -> None:
    draw = ImageDraw.Draw(img)
    top = _header(draw, slide.title, slide.caption, page)
    nodes = slide.nodes[:4]
    area = _content_area(top)
    n = len(nodes)
    if IS_PORTRAIT:
        cols = 1 if n > 2 else min(2, n)
    else:
        cols = 2 if n > 1 else 1
    boxes = _layout_grid(n, area, cols)
    for i, node in enumerate(nodes):
        _draw_node(draw, boxes[i], node)


def render_layers(slide: Slide, img: Image.Image, page: int) -> None:
    draw = ImageDraw.Draw(img)
    top = _header(draw, slide.title, slide.caption, page)
    nodes = slide.nodes[:8]
    area = _content_area(top)
    boxes = _layout_boxes_vertical(len(nodes), area, min_box_h=56)
    for i, node in enumerate(nodes):
        x1, y1, x2, y2 = boxes[i]
        color = PALETTE[i % len(PALETTE)]
        _rounded_box(draw, (x1, y1, x2, y2), fill=(25, 35, 55), outline=color, radius=14, width=2)
        _text_centered(draw, (x1, y1, x2, y2), node.label, node.detail, _font(26), _font(18), C_TITLE)
    for i in range(len(boxes) - 1):
        cx = (boxes[i][0] + boxes[i][2]) // 2
        _arrow(draw, cx, boxes[i][3], cx, boxes[i + 1][1], PALETTE[0])


def render_hub(slide: Slide, img: Image.Image, page: int) -> None:
    draw = ImageDraw.Draw(img)
    top = _header(draw, slide.title, slide.caption, page)
    area = _content_area(top)
    left, ct, right, cb = area
    hub_h = max(72, (cb - ct) // 9)
    cx = (left + right) // 2
    _rounded_box(draw, (left, ct, right, ct + hub_h), fill=(56, 189, 248), outline=PALETTE[0], radius=18, width=3)
    _text_centered(draw, (left, ct, right, ct + hub_h), slide.title[:16], "", _font(28), _font(20), C_TITLE)
    sub_area = (left, ct + hub_h + 14, right, cb)
    satellites = slide.nodes[:5]
    boxes = _layout_boxes_vertical(len(satellites), sub_area, min_box_h=56)
    for i, node in enumerate(satellites):
        _draw_node(draw, boxes[i], node)
        if i == 0:
            _arrow(draw, cx, ct + hub_h, cx, boxes[i][1], PALETTE[0])
        else:
            _arrow(draw, cx, boxes[i - 1][3], cx, boxes[i][1], PALETTE[i % len(PALETTE)])


def render_slide(slide: Slide, page: int, out: Path) -> None:
    img = _gradient_bg()
    kind = slide.kind
    if kind == "title":
        render_title(slide, img, page)
    elif kind in ("flow", "pipeline"):
        render_flow(slide, img, page, vertical=kind == "pipeline")
    elif kind == "cards":
        render_cards(slide, img, page)
    elif kind == "layers":
        render_layers(slide, img, page)
    elif kind == "hub":
        render_hub(slide, img, page)
    else:
        render_cards(slide, img, page)
    img.save(out)


def build_pptx(image_paths: List[Path], out: Path) -> None:
    prs = Presentation()
    if IS_PORTRAIT:
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.333)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in image_paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(p.resolve()), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out)


def build_video_with_narration(
    image_dir: Path,
    audio_dir: Path,
    durations: List[float],
    out: Path,
    min_pad: float = 0.4,
) -> None:
    images = sorted(image_dir.glob("slide_*.png"))
    audios = sorted(audio_dir.glob("slide_*.mp3"))
    if not images:
        raise RuntimeError("无幻灯片图片")
    slide_durations = [max(d + min_pad, 2.5) for d in durations[:len(images)]]
    while len(slide_durations) < len(images):
        slide_durations.append(4.0)

    list_file = image_dir / "ffmpeg_list.txt"
    lines = []
    for img, dur in zip(images, slide_durations):
        lines.append(f"file '{img.resolve()}'")
        lines.append(f"duration {dur:.3f}")
    lines.append(f"file '{images[-1].resolve()}'")
    list_file.write_text("\n".join(lines), encoding="utf-8")

    silent = image_dir / "_silent.mp4"
    subprocess.run(
        [
            "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", str(list_file),
            "-vf", f"fps=30,format=yuv420p,scale={W}:{H}",
            "-c:v", "libx264", "-pix_fmt", "yuv420p", str(silent),
        ],
        check=True,
        capture_output=True,
    )

    if audios:
        audio_list = audio_dir / "audio_concat.txt"
        audio_list.write_text(
            "\n".join(f"file '{a.resolve()}'" for a in audios[:len(images)]),
            encoding="utf-8",
        )
        narration = audio_dir / "_narration.mp3"
        subprocess.run(
            ["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", str(audio_list), "-c", "copy", str(narration)],
            check=True,
            capture_output=True,
        )
        subprocess.run(
            [
                "ffmpeg", "-y", "-i", str(silent), "-i", str(narration),
                "-c:v", "copy", "-c:a", "aac", "-b:a", "192k", "-shortest",
                str(out),
            ],
            check=True,
            capture_output=True,
        )
    else:
        silent.rename(out)


def build_video(image_dir: Path, out: Path, duration: float) -> None:
    images = sorted(image_dir.glob("slide_*.png"))
    list_file = image_dir / "ffmpeg_list.txt"
    lines = []
    for img in images:
        lines.append(f"file '{img.resolve()}'")
        lines.append(f"duration {duration}")
    lines.append(f"file '{images[-1].resolve()}'")
    list_file.write_text("\n".join(lines), encoding="utf-8")
    subprocess.run(
        [
            "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", str(list_file),
            "-vf", "fps=30,format=yuv420p", "-c:v", "libx264", "-pix_fmt", "yuv420p",
            str(out),
        ],
        check=True,
        capture_output=True,
    )


def main():
    parser = argparse.ArgumentParser(description="从 Markdown 生成框图风格 PPT/视频（支持竖屏+讲解）")
    parser.add_argument("article", type=Path, nargs="?", default=None, help="Markdown 文章路径")
    parser.add_argument("-o", "--output-dir", type=Path, default=Path("ppt_output"))
    parser.add_argument("--seconds", type=float, default=4.5, help="无配音时每页秒数")
    parser.add_argument("--landscape", action="store_true", help="横屏 16:9（默认竖屏 9:16）")
    parser.add_argument("--no-voice", action="store_true", help="不生成讲解配音")
    parser.add_argument("--voice", type=str, default=DEFAULT_VOICE, help="单个发音人 ID（与 --voices 互斥时后者优先）")
    parser.add_argument(
        "--voices",
        type=str,
        default="",
        help="多个发音人，逗号分隔；默认按页轮换，如 XiaoxiaoNeural,YunxiNeural",
    )
    parser.add_argument(
        "--voice-preset",
        type=str,
        default="",
        choices=list(VOICE_PRESETS.keys()),
        help="发音人预设组合：female/male/duo/mix/dialect",
    )
    parser.add_argument(
        "--voice-mode",
        type=str,
        default="rotate",
        choices=["rotate", "single", "by_kind"],
        help="多发音人策略：rotate 按页轮换 | single 只用第一个 | by_kind 按版式分配",
    )
    parser.add_argument(
        "--rate",
        type=str,
        default="+0%",
        help="语速，相对百分比，如 +15%%（更快）或 -10%%（更慢）",
    )
    parser.add_argument(
        "--volume",
        type=str,
        default="+0%",
        help="音量，相对百分比，如 +20%% 或 -10%%",
    )
    parser.add_argument(
        "--pitch",
        type=str,
        default="+0Hz",
        help="音调，如 +5Hz 或 -3Hz",
    )
    parser.add_argument(
        "--list-voices",
        action="store_true",
        help="列出可用中文语音后退出",
    )
    args = parser.parse_args()

    if args.list_voices:
        print_voice_catalog()
        return

    set_orientation(not args.landscape)

    if not args.article:
        parser.error("请提供 Markdown 文章路径")

    voices_list = resolve_voices(
        args.voice,
        [v.strip() for v in args.voices.split(",") if v.strip()] if args.voices else None,
        args.voice_preset or None,
    )

    article = args.article.resolve()
    title, _, slides = parse_article(article)
    safe = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", title)[:40].strip("_")
    out_dir = args.output_dir.resolve() / safe
    img_dir = out_dir / "slides"
    audio_dir = out_dir / "audio"
    img_dir.mkdir(parents=True, exist_ok=True)

    paths: List[Path] = []
    for i, slide in enumerate(slides, 1):
        p = img_dir / f"slide_{i:02d}.png"
        render_slide(slide, i, p)
        paths.append(p)

    suffix = "_竖屏" if IS_PORTRAIT else "_横屏"
    pptx_path = out_dir / f"{safe}{suffix}.pptx"
    video_path = out_dir / f"{safe}{suffix}_讲解.mp4"
    build_pptx(paths, pptx_path)

    if args.no_voice:
        build_video(img_dir, video_path, args.seconds)
    else:
        print("正在生成讲解配音…")
        durations = generate_all_narrations(
            slides, audio_dir, voices_list,
            rate=args.rate, volume=args.volume, pitch=args.pitch,
            voice_mode=args.voice_mode,
        )
        build_video_with_narration(img_dir, audio_dir, durations, video_path)

    orient = "竖屏 9:16" if IS_PORTRAIT else "横屏 16:9"
    print(f"标题: {title}")
    print(f"方向: {orient}")
    print(f"幻灯片: {len(slides)} 页")
    print(f"PPTX: {pptx_path}")
    print(f"视频: {video_path}")
    if not args.no_voice:
        print(f"配音: {audio_dir}")
        labels = [VOICE_LABELS.get(v, v) for v in voices_list]
        print(f"发音人({args.voice_mode}): {' | '.join(labels)}")
        print(f"语速: {args.rate}  音量: {args.volume}  音调: {args.pitch}")


if __name__ == "__main__":
    main()
